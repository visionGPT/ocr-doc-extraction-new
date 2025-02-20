# Opticintellect_updated.py
import openai
import io
import os
import base64
import json
import numpy as np
from PIL import Image
import cv2
from ultralytics import YOLO
import fitz
import logging
import asyncio
from concurrent.futures import ThreadPoolExecutor
import time
import copy
from fastapi import Form
from typing import List, Dict, Any

from pymongo import MongoClient
import urllib.parse
import aiohttp
import time
import os
from motor.motor_asyncio import AsyncIOMotorClient

from google.cloud import storage 
from urllib.parse import urlparse

from fastapi import FastAPI, File, UploadFile, HTTPException, Query
from fastapi.responses import PlainTextResponse, HTMLResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from starlette.templating import Jinja2Templates
from starlette.staticfiles import StaticFiles
from pydantic import BaseModel
from fastapi.datastructures import FormData


from docx2pdf import convert
from pptx import Presentation
from pdf2image import convert_from_path
import pandas as pd
import logging 
import traceback
from datetime import datetime
from tabulate import tabulate
import pdfplumber
import camelot
import re

import tempfile
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import pdfkit
from typing import BinaryIO
from tempfile import SpooledTemporaryFile, TemporaryDirectory
from typing import Optional, List, Dict
from wtforms import Form, FileField
from typing import Optional
import os
from templates import template_manager

from google.cloud import storage
from google.oauth2 import service_account
from google.cloud import storage
from google.oauth2 import service_account

from google.generativeai.types import GenerationConfig
from google.generativeai import GenerativeModel, configure
from google.genai.types import (
    FunctionDeclaration,
    GenerateContentConfig,
    GoogleSearch,
    MediaResolution,
    Part,
    Retrieval,
    SafetySetting,
    Tool,
    ToolCodeExecution,
    VertexAISearch,
)


from templates import template_manager
MODEL_ID = "gemini-2.0-flash-001"
configure(api_key="GOOGLE_API_KEY")
genai_client = GenerativeModel(MODEL_ID)

import asyncio
from concurrent.futures import ThreadPoolExecutor
import logging
from typing import List, Tuple, Dict, Any
from PIL import Image
import traceback

class ParallelBatchProcessor:
    def __init__(self, max_workers: int = 8):
        self.max_workers = max_workers
        self.logger = logging.getLogger(__name__)

    async def process_single_page(
        self, 
        image_info: Tuple[Image.Image, List[int], int], 
        openai_api_key: str
    ) -> Dict[str, Any]:
        """
        Process a single page of the document.
        
        Args:
            image_info: Tuple containing (image, page_numbers, page_boundary)
            openai_api_key: OpenAI API key for processing
            
        Returns:
            Dict containing processing results for the page
        """
        try:
            image, page_numbers, page_boundary = image_info
            
            # Detect elements in the image
            result_image, detected_elements = detect(image, page_numbers, page_boundary)
            
            # Save annotated image
            annotated_image_path = save_annotated_image(
                result_image, 
                f"annotated_page_{page_numbers[0]}.png"
            )
            
            # Create chunks from detected elements
            chunks = improved_intelligent_chunking_with_continuity(
                detected_elements,
                SEGMENT_HIERARCHY,
                max_chunk_size=20
            )
            
            # Process each chunk
            processed_chunks = []
            for chunk_idx, chunk in enumerate(chunks):
                # Create image for chunk
                chunk_image = combine_elements_into_image(image, chunk["elements"])
                annotated_chunk_image, _ = detect(chunk_image)
                
                # Add metadata to chunk
                chunk["chunk_index"] = chunk_idx
                chunk["original_image"] = chunk_image
                chunk["annotated_image"] = annotated_chunk_image
                
                # Process the chunk
                processed_chunk = process_chunk(chunk, openai_api_key)
                processed_chunks.append(processed_chunk)
            
            return {
                "pages": page_numbers,
                "processed_chunks": processed_chunks,
                "detected_elements": detected_elements,
                "annotated_image_path": annotated_image_path
            }
            
        except Exception as e:
            self.logger.error(
                f"Error processing page {page_numbers}: {str(e)}\n{traceback.format_exc()}"
            )
            return {
                "pages": page_numbers,
                "error": str(e),
                "traceback": traceback.format_exc()
            }

    async def process_batch(
        self, 
        batch: List[Tuple[Image.Image, List[int], int]], 
        openai_api_key: str
    ) -> List[Dict[str, Any]]:
        """
        Process a batch of pages in parallel.
        
        Args:
            batch: List of image_info tuples
            openai_api_key: OpenAI API key for processing
            
        Returns:
            List of processing results for the batch
        """
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            loop = asyncio.get_event_loop()
            tasks = [
                loop.run_in_executor(
                    executor,
                    lambda x: asyncio.run(self.process_single_page(x, openai_api_key)),
                    image_info
                )
                for image_info in batch
            ]
            return await asyncio.gather(*tasks)

async def parallel_batch_process_document(
    images_with_info: List[Tuple[Image.Image, List[int], int]], 
    openai_api_key: str,
    batch_size: int = 8
) -> List[Dict[str, Any]]:
    """
    Process document pages in parallel batches.
    
    Args:
        images_with_info: List of tuples containing (image, page_numbers, page_boundary)
        openai_api_key: OpenAI API key for processing
        batch_size: Number of pages to process in parallel
        
    Returns:
        List of processing results for all pages
    """
    processor = ParallelBatchProcessor(max_workers=batch_size)
    results = []
    
    # Process pages in batches
    for i in range(0, len(images_with_info), batch_size):
        batch = images_with_info[i:i + batch_size]
        batch_results = await processor.process_batch(batch, openai_api_key)
        results.extend(batch_results)
    
    return results


def initialize_storage_client():
    """Initialize Google Cloud Storage client with credentials"""
    try:
        # Get the current directory
        current_dir = os.path.dirname(os.path.abspath(__file__))
        
        # Path to credentials file
        credentials_path = os.path.join(current_dir, 'credentials.json')
        
        # Check if credentials file exists
        if not os.path.exists(credentials_path):
            raise FileNotFoundError(f"Credentials file not found at: {credentials_path}")
        
        # Create credentials object directly from file
        credentials = service_account.Credentials.from_service_account_file(credentials_path)
        
        # Initialize storage client with credentials
        storage_client = storage.Client(credentials=credentials)
        
        return storage_client
            
    except FileNotFoundError as e:
        print(f"Credentials file error: {str(e)}")
        raise
    except Exception as e:
        print(f"Failed to initialize storage client: {str(e)}")
        raise

# Initialize the storage client
try:
    storage_client = initialize_storage_client()
    bucket_name = "vision-bucket-ai"
    bucket = storage_client.bucket(bucket_name)
    print("Successfully initialized Google Cloud Storage client!")
except Exception as e:
    print(f"Error initializing storage client: {str(e)}")
    raise


# MongoDB setup
username = urllib.parse.quote_plus("sridhargd1234")
password = urllib.parse.quote_plus("charlie@123456")
mongo_client = MongoClient(
    f"mongodb+srv://{username}:{password}@cluster0.f4prp.mongodb.net/?retryWrites=true&w=majority&appName=Cluster0"
)
db = mongo_client["document_processor_db"]
collection = db["processed_documents"]

# Define ProcessResponse model
class ProcessResponse(BaseModel):
    result: str
    chunks: list
    gcs_url: Optional[str] = None  # Make it optional for backward compatibility

# Configure logging
logging.basicConfig(filename='chunk_processing.log', level=logging.INFO, format='%(asctime)s - %(message)s')
# Configure logging
logging.basicConfig(level=logging.ERROR)
logger = logging.getLogger(__name__)


# Define constants with more subtle colors
ENTITIES_COLORS = {
    "Caption": (200, 150, 100),
    "Footnote": (100, 100, 150),
    "Formula": (150, 120, 100),
    "List-item": (180, 200, 150),
    "Page-footer": (100, 120, 150),
    "Page-header": (120, 150, 140),
    "Picture": (220, 150, 160),
    "Section-header": (100, 180, 170),
    "Table": (160, 170, 170),
    "Text": (100, 170, 220),
    "Title": (200, 130, 100),
    "Unknown": (128, 128, 128),
}
BOX_PADDING = 2

# Define the hierarchy for chunking
SEGMENT_HIERARCHY = [
    "Section-header",
    "Title",
    "Page-header",
    "Page-footer",
    "Table",
    "Picture",
    "Caption",
    "Formula",
    "Text",
    "List-item",
    "Footnote",
    "Unknown",
]

# Processing flags
PROCESS_ALL_PAGES = True
PAGE_TO_PROCESS = 2

# Model settings
MODEL_PATH = "models/yolov10x_best.pt"
FILE_ID = "1jTF4xd0Pu7FDFpLTfSGjgTTolZju4_j7"
MODEL_URL = f"https://drive.google.com/uc?id={FILE_ID}"


# Define allowed file types and their MIME types
ALLOWED_EXTENSIONS: Dict[str, str] = {
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'csv': 'text/csv'
}

DEFAULT_TEST_FILES = {
    "invoice": "E:\\ATTACK CAPITAL\\Opticintellect-FastAPI\\default_files\\InvoiceExample.pdf",
    "balance_sheet": "E:\\ATTACK CAPITAL\\Opticintellect-FastAPI\\default_files\\Balance Sheet.pdf",
    "bank_statement": "E:\\ATTACK CAPITAL\\Opticintellect-FastAPI\\default_files\\Bank Statement.pdf"
}

# Create directory for default files if it doesn't exist
if not os.path.exists("default_files"):
    os.makedirs("default_files")

PREBUILT_TEMPLATES = {
    "invoice-template": {
        "template_id": "invoice-template",
        "name": "Invoice Template",
        "description": "Template for processing invoice PDF files",
        "type": "invoice",
        "fields": [
            "invoice_number",
            "date",
            "due_date",
            "company_name",
            "billing_address",
            "line_items",
            "subtotal",
            "tax",
            "total_amount"
        ],
        "sample_file": "templates/invoice_template.pdf"
    },
    "balance-sheet-template": {
        "template_id": "balance-sheet-template",
        "name": "Balance Sheet Template",
        "description": "Template for processing balance sheet PDF files",
        "type": "financial",
        "fields": [
            "reporting_period",
            "assets",
            "liabilities",
            "equity",
            "total_assets",
            "total_liabilities",
            "total_equity"
        ],
        "sample_file": "templates/balance_sheet_template.pdf"
    },
    "bank-statement-template": {
        "template_id": "bank-statement-template",
        "name": "Bank Statement Template",
        "description": "Template for processing bank account statement PDF files",
        "type": "financial",
        "fields": [
            "account_number",
            "statement_period",
            "opening_balance",
            "closing_balance",
            "transactions",
            "deposits",
            "withdrawals",
            "fees"
        ],
        "sample_file": "templates/bank_statement_template.pdf"
    }
}

class TemplateConfig(BaseModel):
    name: str
    description: str
    type: str
    fields: List[str]
    detection_config: dict = {}
    processing_config: dict = {}
    chunk_config: dict = {}
    hierarchy: List[str] = []

class AsyncTemplateManager:
    """Async Template Manager for handling document processing templates"""
    
    def __init__(self):
        self.templates: Dict[str, dict] = PREBUILT_TEMPLATES
        self._lock = asyncio.Lock()
    
    async def get_template(self, template_id: str) -> Optional[dict]:
        """Get a template by ID"""
        async with self._lock:
            return self.templates.get(template_id)
    
    async def add_template(self, template_id: str, template_config: TemplateConfig) -> bool:
        """Add a new template"""
        async with self._lock:
            if template_id in self.templates:
                return False
            self.templates[template_id] = template_config.dict()
            return True
    
    async def update_template(self, template_id: str, template_config: TemplateConfig) -> bool:
        """Update an existing template"""
        async with self._lock:
            if template_id not in self.templates:
                return False
            self.templates[template_id] = template_config.dict()
            return True
    
    async def delete_template(self, template_id: str) -> bool:
        """Delete a template"""
        async with self._lock:
            if template_id not in self.templates:
                return False
            del self.templates[template_id]
            return True
    
    async def list_templates(self) -> List[dict]:
        """List all available templates"""
        async with self._lock:
            return [
                {
                    "template_id": tid,
                    "name": tdata["name"],
                    "description": tdata["description"],
                    "type": tdata["type"]
                }
                for tid, tdata in self.templates.items()
            ]
    
    async def process_template_request(self, file: UploadFile, template_id: str, form_data: FormData) -> dict:
        """
        Process a template-based document request with proper form data handling
        """
        template = await self.get_template(template_id)
        if not template:
            raise ValueError(f"Template '{template_id}' not found")
            
        # Extract form fields based on template configuration
        field_values = {}
        for field in template.get('fields', []):
            if field in form_data:
                field_values[field] = form_data.getlist(field)
            
        return {
            "template": template,
            "file": file,
            "fields": field_values
        }
        
# Initialize template manager
# template_manager = AsyncTemplateManager()

def validate_mime_type(content_type: str) -> str:
    """Validate and return the file extension for the given MIME type."""
    for ext, mime in ALLOWED_EXTENSIONS.items():
        if mime == content_type:
            return ext
    raise HTTPException(
        status_code=400,
        detail=f"Unsupported file type: {content_type}. Supported types: {list(ALLOWED_EXTENSIONS.values())}"
    )

async def convert_docx_to_pdf(input_path: str, output_path: str):
    """Convert DOCX to PDF using docx2pdf."""
    try:
        convert(input_path, output_path)
    except Exception as e:
        logger.error(f"DOCX conversion error: {str(e)}")
        raise Exception(f"DOCX to PDF conversion failed: {str(e)}")

async def convert_spreadsheet_to_pdf(input_path: str, output_path: str, file_type: str):
    """Convert Excel/CSV to PDF using pandas and reportlab."""
    try:
        # Read the spreadsheet
        if file_type == ALLOWED_EXTENSIONS['xlsx']:
            df = pd.read_excel(input_path)
        else:
            df = pd.read_csv(input_path)
        
        # Initialize PDF
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        from reportlab.lib.units import inch
        
        # Create PDF document
        doc = SimpleDocTemplate(
            output_path,
            pagesize=landscape(letter),
            rightMargin=30,
            leftMargin=30,
            topMargin=30,
            bottomMargin=30
        )
        
        # Convert dataframe to list of lists
        data = [df.columns.tolist()]  # Headers
        data.extend(df.values.tolist())
        
        # Create table
        table = Table(data)
        
        # Add style
        style = TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.white),
            ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ])
        
        table.setStyle(style)
        
        # Build PDF
        elements = []
        elements.append(table)
        doc.build(elements)
        
        return True
        
    except Exception as e:
        logger.error(f"Spreadsheet conversion error: {str(e)}\n{traceback.format_exc()}")
        raise Exception(f"Spreadsheet to PDF conversion failed: {str(e)}")

async def convert_pptx_to_pdf(input_path: str, output_path: str):
    """
    Convert PPTX to PDF using python-pptx and reportlab with improved error handling
    and better text/shape rendering.
    """
    try:
        # Load the presentation
        prs = Presentation(input_path)
        
        # Initialize PDF canvas with letter size
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        # Create a temporary directory for slide processing
        with TemporaryDirectory() as temp_dir:
            for i, slide in enumerate(prs.slides):
                # Initialize y_offset for the current slide
                y_offset = 50  # Start from top with margin
                
                # Process each shape in the slide
                for shape in slide.shapes:
                    # Handle text shapes
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        
                        if hasattr(shape, 'name') and 'Title' in shape.name:
                            # Title text
                            c.setFont("Helvetica-Bold", 24)
                            c.drawString(50, height - y_offset, text)
                            y_offset += 40
                        else:
                            # Regular text
                            c.setFont("Helvetica", 12)
                            text_object = c.beginText(50, height - y_offset)
                            
                            # Handle text wrapping
                            wrapped_text = "\n".join(
                                [line.strip() for line in text.split('\n') if line.strip()]
                            )
                            
                            for line in wrapped_text.split('\n'):
                                if line.strip():
                                    text_object.textLine(line.strip())
                                    y_offset += 20
                            
                            c.drawText(text_object)
                    
                    # Handle image shapes
                    elif hasattr(shape, 'image'):
                        try:
                            # Extract and process image
                            image_stream = io.BytesIO(shape.image.blob)
                            img = Image.open(image_stream)
                            
                            # Calculate image dimensions
                            img_width, img_height = img.size
                            scale = min(
                                (width - 100) / img_width,
                                (height - y_offset - 50) / img_height,
                                1.0
                            )
                            
                            new_width = int(img_width * scale)
                            new_height = int(img_height * scale)
                            
                            # Save and add image to PDF
                            temp_img_path = os.path.join(temp_dir, f"img_{i}.png")
                            img.save(temp_img_path)
                            
                            c.drawImage(
                                temp_img_path,
                                50,
                                height - y_offset - new_height,
                                width=new_width,
                                height=new_height
                            )
                            
                            y_offset += new_height + 20
                            
                        except Exception as img_error:
                            logger.warning(f"Failed to process image in slide {i}: {str(img_error)}")
                            continue
                
                    # Check if new page is needed
                    if y_offset > height - 100:
                        c.showPage()
                        y_offset = 50
                
                # Add page number at the bottom
                c.setFont("Helvetica", 10)
                c.drawString(width - 50, 30, f"Page {i + 1}")
                
                # Move to next page for next slide
                c.showPage()
            
            # Save and verify the final PDF
            c.save()
            
            if not os.path.exists(output_path):
                raise Exception("PDF file was not created successfully")
            
            if os.path.getsize(output_path) == 0:
                raise Exception("Created PDF file is empty")
            
            return True
            
    except Exception as e:
        logger.error(f"PPTX conversion error: {str(e)}\n{traceback.format_exc()}")
        raise Exception(f"PPTX to PDF conversion failed: {str(e)}")

async def convert_to_pdf(file_content: bytes, file_type: str) -> bytes:
    """Main conversion function that handles all supported file types."""
    with TemporaryDirectory() as temp_dir:
        try:
            # Get file extension from MIME type
            file_ext = validate_mime_type(file_type)
            
            # Create temporary file paths
            temp_input = os.path.join(temp_dir, f"input.{file_ext}")
            temp_output = os.path.join(temp_dir, "output.pdf")
            
            # Write input file
            with open(temp_input, 'wb') as f:
                f.write(file_content)
            
            # Convert based on file type
            if file_type == ALLOWED_EXTENSIONS['docx']:
                await convert_docx_to_pdf(temp_input, temp_output)
            elif file_type in [ALLOWED_EXTENSIONS['xlsx'], ALLOWED_EXTENSIONS['csv']]:
                await convert_spreadsheet_to_pdf(temp_input, temp_output, file_type)
            elif file_type == ALLOWED_EXTENSIONS['pptx']:
                success = await convert_pptx_to_pdf(temp_input, temp_output)
                if not success:
                    raise Exception("PPTX conversion failed without specific error")
            
            # Read and return the PDF content
            if os.path.exists(temp_output):
                with open(temp_output, 'rb') as f:
                    pdf_content = f.read()
                    if len(pdf_content) == 0:
                        raise Exception("Generated PDF is empty")
                    return pdf_content
            else:
                raise Exception("PDF output file was not created")
                
        except Exception as e:
            logger.error(f"Conversion error: {str(e)}\n{traceback.format_exc()}")
            raise HTTPException(
                status_code=500,
                detail=f"Document conversion failed: {str(e)}"
            )


DEFAULT_TEST_FILES = {
    "invoice": "E:\\ATTACK CAPITAL\\Opticintellect-FastAPI\\default_files\\InvoiceExample.pdf",
    "balance_sheet": "E:\\ATTACK CAPITAL\\Opticintellect-FastAPI\\default_files\\Balance Sheet.pdf",
    "bank_statement": "E:\\ATTACK CAPITAL\\Opticintellect-FastAPI\\default_files\\Bank Statement.pdf"
}   

# Add Jinja2 Templates for rendering
templates = Jinja2Templates(directory="templates")

# Initialize OpenAI API key configuration
class OpenAIConfig(BaseModel):
    api_key: str

# Global variables
openai_config = None
DETECTION_MODEL = None

# Configuration model for Google API
class GoogleAPIConfig(BaseModel):
    api_key: str

# Global configuration variable
google_api_config = None

def save_annotated_image(image, filename="annotated_image.png"):
    """
    Save the annotated image to a static directory
    """
    if not os.path.exists("static"):
        os.makedirs("static")
    
    image.save(os.path.join("static", filename))
    return f"/static/{filename}"

def draw_box_and_label(image, start_box, end_box, cls, detection_class_conf):
    box_scale_factor = 0.001
    label_scale_factor = 0.5

    line_thickness = max(
        round(box_scale_factor * (image.shape[0] + image.shape[1]) / 2), 1
    )

    cv2.rectangle(
        img=image,
        pt1=start_box,
        pt2=end_box,
        color=ENTITIES_COLORS.get(cls, (128, 128, 128)),
        thickness=line_thickness,
    )

    text = f"{cls} {detection_class_conf:.2f}"
    font_scale = label_scale_factor
    font_thickness = max(line_thickness - 1, 1)

    (text_w, text_h), _ = cv2.getTextSize(
        text, cv2.FONT_HERSHEY_SIMPLEX, fontScale=font_scale, thickness=font_thickness
    )

    cv2.rectangle(
        image,
        (start_box[0], start_box[1] - text_h - BOX_PADDING * 2),
        (start_box[0] + text_w + BOX_PADDING * 2, start_box[1]),
        ENTITIES_COLORS.get(cls, (128, 128, 128)),
        thickness=-1,
    )

    cv2.putText(
        image,
        text,
        (start_box[0] + BOX_PADDING, start_box[1] - BOX_PADDING),
        cv2.FONT_HERSHEY_SIMPLEX,
        fontScale=font_scale,
        color=(255, 255, 255),
        thickness=font_thickness,
    )

def load_model():
    global DETECTION_MODEL
    if not os.path.exists(MODEL_PATH):
        os.makedirs(os.path.dirname(MODEL_PATH), exist_ok=True)
        try:
            import gdown
            gdown.download(MODEL_URL, MODEL_PATH, quiet=False)
        except Exception as e:
            logging.error(f"Error downloading the model: {e}")
            raise HTTPException(status_code=500, detail=f"Error downloading model: {e}")
    
    DETECTION_MODEL = YOLO(MODEL_PATH)
    return DETECTION_MODEL

def detect(image, page_numbers=None, page_boundary=None):
    image_cv = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
    results = DETECTION_MODEL.predict(source=image_cv, conf=0.2, iou=0.8)
    boxes = results[0].boxes

    detected_elements = []

    if len(boxes) == 0:
        return Image.fromarray(cv2.cvtColor(image_cv, cv2.COLOR_BGR2RGB)), detected_elements

    for box in boxes:
        detection_class_conf = box.conf.item()
        cls_index = int(box.cls)
        cls = list(ENTITIES_COLORS.keys())[cls_index] if cls_index < len(ENTITIES_COLORS) else "Unknown"

        start_box = (int(box.xyxy[0][0]), int(box.xyxy[0][1]))
        end_box = (int(box.xyxy[0][2]), int(box.xyxy[0][3]))

        if page_boundary is not None and page_numbers is not None:
            if end_box[1] <= page_boundary:
                page_number = page_numbers[0]
            elif start_box[1] >= page_boundary:
                page_number = page_numbers[1]
            else:
                page_number = page_numbers
        else:
            page_number = page_numbers[0] if page_numbers else None

        detected_elements.append({
            "class": cls,
            "confidence": detection_class_conf,
            "coordinates": {"start": start_box, "end": end_box},
            "page_number": page_number,
        })

        draw_box_and_label(image_cv, start_box, end_box, cls, detection_class_conf)

    detected_elements.sort(key=lambda x: (x["coordinates"]["start"][1], x["coordinates"]["start"][0]))
    for idx, element in enumerate(detected_elements):
        element["index"] = idx

    return Image.fromarray(cv2.cvtColor(image_cv, cv2.COLOR_BGR2RGB)), detected_elements

def format_table_as_markdown(table_data):
    """
    Convert table data to well-structured markdown format with proper alignment and spacing.
    """
    if not table_data or 'raw_data' not in table_data:
        return "Error: Invalid table data"

    header = table_data['raw_data']['header']
    data = table_data['raw_data']['data']

    # Clean and standardize data
    header = [str(h).strip() if h is not None else '' for h in header]
    cleaned_data = []
    for row in data:
        cleaned_row = [str(cell).strip() if cell is not None else '' for cell in row]
        if any(cleaned_row):  # Skip entirely empty rows
            cleaned_data.append(cleaned_row)

    # Calculate column widths
    col_widths = [len(h) for h in header]
    for row in cleaned_data:
        for i, cell in enumerate(row):
            col_widths[i] = max(col_widths[i], len(cell))
    
    # Add padding to column widths
    col_widths = [w + 2 for w in col_widths]  # Add 2 spaces padding per column

    # Create header row with proper spacing
    header_cells = [f" {h.ljust(w-2)} " for h, w in zip(header, col_widths)]
    md_table = "|" + "|".join(header_cells) + "|\n"
    
    # Create separator row with proper alignment indicators
    separator_cells = ['-' * w for w in col_widths]
    md_table += "|" + "|".join(separator_cells) + "|\n"
    
    # Create data rows with proper spacing
    for row in cleaned_data:
        # Ensure row has same number of columns as header
        while len(row) < len(header):
            row.append('')
        
        # Format each cell with proper spacing
        formatted_cells = [f" {cell.ljust(w-2)} " for cell, w in zip(row, col_widths)]
        md_table += "|" + "|".join(formatted_cells) + "|\n"
    
    return md_table


def extract_tables_with_pdfplumber(pdf_content):
    """
    Extract tables from PDF using pdfplumber with improved markdown formatting.
    """
    tables = []
    try:
        with pdfplumber.open(io.BytesIO(pdf_content)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_tables = page.extract_tables()
                for table_num, table in enumerate(page_tables, 1):
                    if table:
                        # Clean and standardize table data
                        header = [str(cell).strip() if cell is not None else '' for cell in table[0]]
                        data = []
                        for row in table[1:]:
                            cleaned_row = [str(cell).strip() if cell is not None else '' for cell in row]
                            if any(cleaned_row):  # Skip entirely empty rows
                                data.append(cleaned_row)

                        if header and data:  # Only process non-empty tables
                            table_data = {
                                'raw_data': {
                                    'header': header,
                                    'data': data
                                }
                            }
                            
                            # Convert to properly formatted markdown
                            markdown_table = format_table_as_markdown(table_data)
                            
                            tables.append({
                                'page': page_num,
                                'table_number': table_num,
                                'content': markdown_table,
                                'raw_data': table_data['raw_data']
                            })
    except Exception as e:
        logging.error(f"Error extracting tables with pdfplumber: {str(e)}")
        return []

    return tables


def extract_tables_with_camelot(pdf_path):
    """
    Extract tables from PDF using Camelot with improved markdown formatting.
    """
    tables = []
    try:
        camelot_tables = camelot.read_pdf(pdf_path, pages='all', flavor='stream')
        for i, table in enumerate(camelot_tables):
            df = table.df
            
            # Clean and standardize header and data
            header = [str(col).strip() for col in df.columns]
            data = []
            for _, row in df.iterrows():
                cleaned_row = [str(cell).strip() for cell in row]
                if any(cleaned_row):  # Skip entirely empty rows
                    data.append(cleaned_row)

            if header and data:  # Only process non-empty tables
                table_data = {
                    'raw_data': {
                        'header': header,
                        'data': data
                    }
                }
                
                # Convert to properly formatted markdown
                markdown_table = format_table_as_markdown(table_data)
                
                tables.append({
                    'page': table.page,
                    'table_number': i + 1,
                    'content': markdown_table,
                    'raw_data': table_data['raw_data']
                })
    except Exception as e:
        logging.error(f"Error extracting tables with Camelot: {str(e)}")
        return []

    return tables

def clean_table_formatting(table_content):
    """
    Clean table formatting by replacing special characters with dashes.
    """
    # Replace plus signs and equals signs with dashes, preserving structure
    cleaned_content = re.sub(r'[+=]', '-', table_content)
    return cleaned_content

def process_chunk(chunk, openai_api_key):
    """
    Enhanced chunk processing without summary generation.
    """
    MAX_RETRIES = 3
    BASE_DELAY = 2  # Base delay in seconds
    
    try:
        annotated_image = chunk["annotated_image"]
        chunk_classes = [element["class"] for element in chunk["elements"]]
        dominant_class = determine_dominant_class(chunk_classes)
        
        # Get enhanced prompts
        system_instruction, prompt = get_prompt_for_class(dominant_class)
        
        # Convert image to base64
        img_str = encode_image(annotated_image)
        
        # Create enhanced Gemini prompt parts
        prompt_parts = [
            {"text": system_instruction},
            {"text": prompt},
            {"text": "Important: Extract ALL content completely without any omissions."},
            {"inline_data": {
                "mime_type": "image/png",
                "data": img_str
            }}
        ]

        # Initialize variables for retry loop
        retry_count = 0
        last_error = None
        
        while retry_count < MAX_RETRIES:
            try:
                # Configure Gemini with optimized parameters
                generation_config = GenerationConfig(
                    temperature=0.1,
                    top_p=0.99,
                    top_k=40,
                    max_output_tokens=8000,
                    candidate_count=1
                )
                
                # Process with Gemini
                response = genai_client.generate_content(
                    prompt_parts,
                    generation_config=generation_config,
                    safety_settings=[
                        {
                            "category": "HARM_CATEGORY_DANGEROUS",
                            "threshold": "BLOCK_NONE"
                        }
                    ]
                )
                
                # Validate and ensure complete extraction
                extracted_content = response.text.strip()
                
                # If content seems truncated, try extracting again with different parameters
                if len(extracted_content.split()) < 10 or extracted_content.endswith('...'):
                    generation_config.temperature = 0.2
                    generation_config.top_p = 1.0
                    response = genai_client.generate_content(
                        prompt_parts,
                        generation_config=generation_config
                    )
                    extracted_content = response.text.strip()
                
                return {
                    "chunk_index": chunk["chunk_index"],
                    "class": dominant_class,
                    "result": extracted_content,
                    "elements": chunk["elements"],
                }
                
            except Exception as e:
                last_error = e
                retry_count += 1
                
                if "429" in str(e) or "quota" in str(e).lower():
                    delay = BASE_DELAY * (2 ** (retry_count - 1))
                    logging.warning(f"Rate limit hit, retrying in {delay} seconds...")
                    time.sleep(delay)
                    continue
                    
                logging.error(f"Chunk processing error (attempt {retry_count}): {str(e)}")
                if retry_count < MAX_RETRIES:
                    time.sleep(BASE_DELAY)
                    continue
                else:
                    break
        
        error_msg = f"Error after {MAX_RETRIES} retries: {str(last_error)}"
        logging.error(error_msg)
        return {
            "chunk_index": chunk.get("chunk_index", -1),
            "class": dominant_class,
            "result": error_msg,
            "elements": chunk.get("elements", []),
        }
        
    except Exception as e:
        logging.error(f"Unexpected error in chunk processing: {str(e)}\n{traceback.format_exc()}")
        return {
            "chunk_index": chunk.get("chunk_index", -1),
            "class": dominant_class if 'dominant_class' in locals() else "Unknown",
            "result": f"Unexpected error: {str(e)}",
            "elements": chunk.get("elements", []),
        }


def process_pdf(pdf_document):
    """
    Process a PDF document and return a list of combined images from two pages.
    """
    images = []
    total_pages = len(pdf_document)

    if PROCESS_ALL_PAGES:
        pages_to_process = range(0, total_pages, 2)  # Process two pages at a time
    else:
        # Ensure PAGE_TO_PROCESS is within the valid range
        if PAGE_TO_PROCESS < 0 or PAGE_TO_PROCESS >= total_pages:
            logging.error(f"PAGE_TO_PROCESS {PAGE_TO_PROCESS} is out of range.")
            return images
        pages_to_process = [PAGE_TO_PROCESS]

    for i in pages_to_process:
        # Load first page
        page1 = pdf_document.load_page(i)
        pix1 = page1.get_pixmap(matrix=fitz.Matrix(2, 2))  # Increase resolution
        img1 = Image.frombytes("RGB", [pix1.width, pix1.height], pix1.samples)

        if i + 1 < total_pages:
            # Load second page
            page2 = pdf_document.load_page(i + 1)
            pix2 = page2.get_pixmap(matrix=fitz.Matrix(2, 2))
            img2 = Image.frombytes("RGB", [pix2.width, pix2.height], pix2.samples)

            # Combine images vertically
            combined_image = Image.new("RGB", (max(img1.width, img2.width), img1.height + img2.height))
            combined_image.paste(img1, (0, 0))
            combined_image.paste(img2, (0, img1.height))

            images.append((combined_image, [i, i + 1], img1.height))
        else:
            images.append((img1, [i], img1.height))

    return images

def combine_elements_into_image(image, elements):
    """
    Combines detected elements into an image by cropping and stitching.
    """
    if not elements:
        return image

    # Proceed with the existing logic
    x_start = min(element["coordinates"]["start"][0] for element in elements)
    x_end = max(element["coordinates"]["end"][0] for element in elements)
    y_start = min(element["coordinates"]["start"][1] for element in elements)
    y_end = max(element["coordinates"]["end"][1] for element in elements)
    
    # Crop the image to the bounding box of all elements
    cropped_image = image.crop((x_start, y_start, x_end, y_end))
    
    return cropped_image

def determine_dominant_class(classes):
    """
    Determine the dominant class in a list of classes based on hierarchy.
    """
    for cls in SEGMENT_HIERARCHY:
        if cls in classes:
            return cls
    return "Unknown"

def encode_image(image):
    """
    Encode PIL image to base64 string.
    """
    # Maintain aspect ratio while resizing
    max_size = (800, 800)
    image.thumbnail(max_size, Image.LANCZOS)

    buffered = io.BytesIO()
    image.save(buffered, format="PNG")
    img_str = base64.b64encode(buffered.getvalue()).decode()

    return img_str

def get_prompt_for_class(cls):
    """
    Get the system instruction and prompt for a given class.
    Returns a tuple of (system_instruction, prompt)
    """
    prompts = {
        "Text": (
            """You are a precise text extraction assistant. Your task is to extract ALL text content completely and accurately, without omitting any information. Always ensure you capture the full context and complete paragraphs.""",
            """Please extract ALL text content from this image with these requirements:
            1. Extract every single word and sentence completely
            2. Maintain original paragraph structure and formatting
            3. Do not summarize or truncate content
            4. Include all contextual information
            5. Extract content exactly as shown, without any omissions
            
            Return the complete extracted text content:"""
        ),
        "Section-header": (
            """You are a section header extraction specialist. Extract complete headers with their full context.""",
            """Extract ALL section headers from this image with these requirements:
            1. Include complete header text
            2. Maintain any subheaders or related context
            3. Preserve formatting and hierarchy
            4. Do not truncate or summarize
            
            Return the complete headers:"""
        ),
        "Skills": (
            """You are a skills extraction specialist. Extract ALL skills mentioned completely.""",
            """Extract ALL skills from this image with these requirements:
            1. Include every single skill mentioned
            2. Include both technical and soft skills
            3. Maintain original grouping/categorization
            4. Do not omit any skills
            5. Extract complete skill descriptions if present
            
            Return the complete list of skills:"""
        ),
        "Title": (
            "You are a title extraction expert using Gemini AI. Extract titles with precise formatting.",
            "Extract and return only the title text from this image, maintaining original emphasis and style."
        ),
        "Caption": (
            "You are a caption extraction specialist using Gemini AI. Extract captions with context.",
            "Extract and return only the caption text from this image, preserving its relationship to content."
        ),
        "Footnote": (
            "You are a footnote extraction expert using Gemini AI. Extract footnotes with references.",
            "Extract and return only the footnote text from this image, including reference markers."
        ),
        "Page-header": (
            "You are a page header extraction specialist using Gemini AI. Extract headers precisely.",
            "Extract and return only the page header text from this image, maintaining formatting."
        ),
        "Page-footer": (
            "You are a page footer extraction specialist using Gemini AI. Extract footers exactly.",
            "Extract and return only the page footer text from this image, preserving format and structure."
        ),
        "List-item": (
            "You are a list extraction expert using Gemini AI. Extract lists with structure.",
            "Extract and return only the list item text from this image, maintaining bullets/numbers and hierarchy."
        ),
        "Table": (
            """You are a table extraction specialist focused on creating perfectly formatted markdown tables.""",
            """Extract and format ALL tables from this image as markdown tables, following these strict requirements:
            1. Use precise markdown table syntax with proper cell alignment
            2. Include table headers with correct formatting
            3. Use consistent column widths based on content
            4. Add proper spacing within cells
            5. Ensure clean data presentation
            
            Required markdown table format:

            | Header 1    | Header 2    | Header 3    |
            |------------|-------------|-------------|
            | Data 1     | Data 2      | Data 3      |
            | Long Data  | Short Data  | Medium Data |

            Extract ALL tables following this exact format. Maintain proper alignment and spacing.
            Do not truncate or omit any content."""
        ),
        "Picture": (
            "You are an image content specialist using Gemini AI. Extract visual content accurately.",
            "Extract and describe any text or data present in this image, focusing on accuracy and completeness."
        ),
        "Formula": (
            "You are a formula extraction expert using Gemini AI. Convert to LaTeX when possible.",
            "Extract and convert the mathematical formula to LaTeX format if possible. Return only the formula."
        ),
        "Unknown": (
            "You are a general content extraction specialist using Gemini AI. Extract all content types.",
            "Extract any recognizable content from this image with high accuracy. Return only the extracted content."
        )
    }
    return prompts.get(cls, (
        "You are a complete content extraction specialist. Extract ALL content without omissions.",
        "Extract ALL content from this image completely and accurately. Do not truncate or summarize any information."
    ))

def generate_overall_summary(results, openai_api_key, max_retries=3, base_delay=1):
    """
    Generate an overall summary for the entire document
    """
    try:
        # Combine all extracted content
        all_content = ""
        for page in results:
            for chunk in page['processed_chunks']:
                if chunk.get('result'):
                    all_content += chunk['result'] + "\n\n"

        client = openai.OpenAI(api_key=openai_api_key)
        
        for attempt in range(max_retries):
            try:
                response = client.chat.completions.create(
                    model="gpt-4o-mini", 
                    messages=[
                        {
                            "role": "system",
                            "content": """You are a comprehensive document summarization expert. Create a short
                            summary capturing key points, main themes, and important details from the document. 
                            Structure your summary with clear sections and highlight critical information."""
                        },
                        {
                            "role": "user",
                            "content": f"Please provide a detailed summary of this document:\n\n{all_content}"
                        }
                    ],
                    max_tokens=1000,  
                    temperature=0.3,
                )
                
                return response.choices[0].message.content.strip()
                
            except Exception as e:
                if "rate_limit" in str(e).lower() and attempt < max_retries - 1:
                    delay = base_delay * (2 ** attempt)
                    time.sleep(delay)
                    continue
                raise e
                
    except Exception as e:
        logging.error(f"Overall summary generation error: {str(e)}")
        return f"Error generating overall summary: {str(e)}"

def improved_intelligent_chunking_with_continuity(detected_elements, hierarchy, max_chunk_size=5):
    """
    Optimized chunking strategy with increased chunk size for better performance
    """
    chunks = []
    current_chunk_elements = []
    priorities = {element_type: idx for idx, element_type in enumerate(hierarchy)}
    current_priority = priorities.get(detected_elements[0]["class"], len(hierarchy)) if detected_elements else None

    def should_start_new_chunk(element, current_priority):
        element_priority = priorities.get(element["class"], len(hierarchy))
        
        # Skip chunking for elements that don't need summaries
        if element["class"] in [
            "Section-header", "Title", "Caption", "Footnote", 
            "Page-header", "Page-footer", "Unknown"
        ]:
            return True
            
        return (
            len(current_chunk_elements) >= max_chunk_size or
            element["class"] in ["Section-header", "Title", "Skills"] or
            (current_priority is not None and abs(element_priority - current_priority) > 3)
        )

    for element in detected_elements:
        if should_start_new_chunk(element, current_priority):
            if current_chunk_elements:
                chunks.append({"elements": current_chunk_elements})
                current_chunk_elements = []
            current_priority = priorities.get(element["class"], len(hierarchy))
        
        current_chunk_elements.append(element)

    # Add remaining elements
    if current_chunk_elements:
        chunks.append({"elements": current_chunk_elements})

    return chunks

def batch_process_chunks(chunks, openai_api_key, batch_size=64):
    """
    Process chunks in larger parallel batches for improved performance
    """
    async def process_batch(batch):
        with ThreadPoolExecutor(max_workers=min(32, len(batch))) as executor:
            loop = asyncio.get_event_loop()
            tasks = [
                loop.run_in_executor(executor, process_chunk, chunk, openai_api_key)
                for chunk in batch
            ]
            return await asyncio.gather(*tasks)
    
    results = []
    for i in range(0, len(chunks), batch_size):
        batch = chunks[i:i + batch_size]
        batch_results = asyncio.run(process_batch(batch))
        results.extend(batch_results)
    
    return results

def format_results_as_markdown(results):
    """
    Convert processing results to markdown format with optimized summary display
    """
    markdown_output = "# Document Extraction Results\n\n"
    
    for page_index, page_result in enumerate(results, 1):
        markdown_output += f"## Page {page_index}\n\n"
        
        # Detected Elements Summary
        markdown_output += "### Detected Elements\n"
        for element in page_result['detected_elements']:
            markdown_output += f"- {element['class']} (Confidence: {element['confidence']:.2f})\n"
        
        markdown_output += "\n### Processed Chunks\n\n"
        
        for chunk_index, chunk in enumerate(page_result['processed_chunks'], 1):
            markdown_output += f"#### Chunk {chunk_index} (Type: {chunk['class']})\n\n"
            markdown_output += f"**Extracted Content:**\n```\n{chunk['result']}\n```\n\n"
            
            # Only display summary if it exists and is not empty
            if chunk.get('summary'):
                markdown_output += f"**Summary:**\n{chunk['summary']}\n\n"
            
            markdown_output += "---\n\n"
    
    return markdown_output

def initialize_gemini(api_key: str):
    """Initialize Gemini with improved configuration."""
    try:
        global genai_client
        configure(api_key=api_key)
        
        # Create model with specific configuration
        genai_client = GenerativeModel(
            MODEL_ID,
            generation_config=GenerationConfig(
                temperature=0.1,
                top_p=0.95,
                top_k=20,
                max_output_tokens=4000,
                candidate_count=1
            )
        )
        
        # Test the model with a simple prompt
        test_response = genai_client.generate_content("Test connection")
        if test_response and hasattr(test_response, 'text'):
            logging.info("Gemini model initialized successfully")
            return True
        else:
            logging.error("Gemini model initialization failed: Invalid response format")
            return False
            
    except Exception as e:
        logging.error(f"Failed to initialize Gemini: {str(e)}")
        return False

# Modify the app to include CORS and static files
app = FastAPI(title="Document Extraction API")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize template manager
template_manager = AsyncTemplateManager()

# Modify the upload page route to include filename input
@app.get("/", response_class=HTMLResponse)
async def upload_page():
    """
    Render the document upload page with filename input
    """
    return """
    <!DOCTYPE html>
    <html>
    <head>
        <title>Document Extraction</title>
    </head>F
    <body>
        <h1>Document Extraction</h1>
        <form action="/upload-document" method="post" enctype="multipart/form-data">
            <div>
                <label for="file">Upload File:</label>
                <input type="file" name="file" id="file" accept=".pdf,.png,.jpg,.jpeg,.bmp">
            </div>
            <div style="margin-top: 10px;">
                <label for="converted_filename">Or Enter Filename from Previous Conversion:</label>
                <input type="text" name="converted_filename" id="converted_filename" placeholder="Enter converted filename">
            </div>
            <div style="margin-top: 10px;">
                <input type="submit" value="Upload and Process">
            </div>
        </form>
    </body>
    </html>
    """

@app.post("/initialize-openai")
async def initialize_openai(config: OpenAIConfig):
    """
    Endpoint to set OpenAI API configuration
    """
    global openai_config
    openai_config = config
    return {"message": "OpenAI configuration initialized successfully"}

@app.post("/initialize-google-api")
async def initialize_google_api(config: GoogleAPIConfig):
    """
    Endpoint to set Google API configuration and initialize Gemini model
    """
    try:
        global google_api_config
        google_api_config = config
        
        if initialize_gemini(config.api_key):
            return {
                "status": "success",
                "message": "Google API configuration initialized successfully",
                "model": MODEL_ID
            }
        else:
            raise HTTPException(
                status_code=500,
                detail="Failed to initialize Gemini model"
            )
        
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to initialize Google API configuration: {str(e)}"
        )
    
# Helper function to check if API is configured
def check_api_configured():
    """Check if the Google API is configured"""
    if not google_api_config:
        raise HTTPException(
            status_code=400,
            detail="Google API key not initialized. Please call /initialize-google-api first."
        )
    return True


# Create a directory for storing converted files
CONVERTED_FILES_DIR = "converted_files"
os.makedirs(CONVERTED_FILES_DIR, exist_ok=True)

@app.post("/process-document", response_class=HTMLResponse)
async def process_document(
    file: UploadFile = File(...),
):
    """
    Unified endpoint to handle document processing.
    Supports both direct PDF uploads and convertible document formats (.docx, .csv, .xlsx, .pptx).
    Returns HTML response with extraction results and file_id.
    """
    try:
        # Check if OpenAI configuration is set
        if not openai_config:
            raise HTTPException(status_code=400, detail="OpenAI API key not initialized")

        # Read file content
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="Empty file provided")

        # Generate a unique file_id early for tracking
        file_id = f"{int(time.time())}_{file.filename}"

        # Upload file to GCS
        file_name = f"{int(time.time())}_{file.filename}"
        blob = bucket.blob(file_name)
        blob.upload_from_string(content)
        gcs_url = f"https://storage.googleapis.com/{bucket_name}/{file_name}"

        # Create initial MongoDB document
        initial_doc = {
            "file_id": file_id,
            "gcs_url": gcs_url,
            "status": "processing",
            "result": "Processing started",
            "chunks": [],
            "file_name": file.filename,
            "created_at": datetime.utcnow()
        }
        collection.insert_one(initial_doc)

        # Determine file type and process accordingly
        file_type = file.content_type
        
        # Handle convertible documents
        if file_type in ALLOWED_EXTENSIONS.values():
            try:
                # Convert to PDF
                pdf_content = await convert_to_pdf(content, file_type)
                
                # Generate unique filename for converted PDF
                converted_filename = f"converted_{os.path.splitext(file.filename)[0]}_{int(time.time())}.pdf"
                converted_filepath = os.path.join(CONVERTED_FILES_DIR, converted_filename)
                
                # Save converted PDF
                with open(converted_filepath, 'wb') as f:
                    f.write(pdf_content)
                
                # Update MongoDB with conversion status
                collection.update_one(
                    {"file_id": file_id},
                    {"$set": {
                        "status": "converted",
                        "converted_filename": converted_filename
                    }}
                )
                
                # Use the converted PDF content for processing
                content = pdf_content
                file_type = "application/pdf"
                display_filename = converted_filename
            except Exception as conv_error:
                collection.update_one(
                    {"file_id": file_id},
                    {"$set": {
                        "status": "conversion_failed",
                        "error": str(conv_error)
                    }}
                )
                raise HTTPException(status_code=500, detail=f"File conversion failed: {str(conv_error)}")
        elif file_type == "application/pdf":
            display_filename = file.filename
        else:
            collection.update_one(
                {"file_id": file_id},
                {"$set": {
                    "status": "invalid_file_type",
                    "error": f"Unsupported file type: {file_type}"
                }}
            )
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {file_type}. Supported types: {list(ALLOWED_EXTENSIONS.values()) + ['application/pdf']}"
            )

        # Process the PDF content
        pdf_document = fitz.open(stream=content, filetype="pdf")
        images_with_info = process_pdf(pdf_document)

        # Update MongoDB with processing started
        collection.update_one(
            {"file_id": file_id},
            {"$set": {
                "status": "processing_pages",
                "total_pages": len(images_with_info)
            }}
        )

        # Process images and collect results
        try:
            results = await parallel_batch_process_document(images_with_info, openai_config.api_key)

            # Collect annotated image paths
            annotated_image_paths = [
                result['annotated_image_path'] 
                for result in results 
                if 'annotated_image_path' in result
            ]

            # Update MongoDB with progress
            for idx, result in enumerate(results):
                if 'error' in result:
                    collection.update_one(
                        {"file_id": file_id},
                        {"$set": {
                            f"page_{idx}_status": "failed",
                            f"page_{idx}_error": result['error']
                        }}
                    )
                else:
                    collection.update_one(
                        {"file_id": file_id},
                        {"$set": {
                            f"page_{idx}_status": "completed"
                        }}
                    )
                    
        except Exception as e:
            logger.error(f"Error in parallel processing: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Error in parallel processing: {str(e)}")
        
        overall_summary = generate_overall_summary(results, openai_config.api_key)

        # Update final results in MongoDB
        final_update = {
            "status": "completed",
            "result": "Processing complete",
            "chunks": results,
            "overall_summary": overall_summary,
            "completed_at": datetime.utcnow(),
            "annotated_image_paths": annotated_image_paths,
            "total_chunks_processed": sum(len(page["processed_chunks"]) for page in results)
        }
        collection.update_one(
            {"file_id": file_id},
            {"$set": final_update}
        )

        # Clean up temporary files
        try:
            if 'converted_filepath' in locals():
                os.remove(converted_filepath)
        except Exception as cleanup_error:
            logger.error(f"Error cleaning up temporary files: {str(cleanup_error)}")

        # Return the HTML response with file_id included
        return templates.TemplateResponse(
            "results.html",
            {
                "request": {"type": "http.request", "method": "POST"},
                "file_name": display_filename,
                "results": results,
                "overall_summary": overall_summary,
                "annotated_image_paths": annotated_image_paths,
                "file_id": file_id,
                "gcs_url": gcs_url,
                "message": f"File ID for this document: {file_id}"
            }
        )

    except HTTPException as he:
        # Update MongoDB with error status for HTTP exceptions
        if 'file_id' in locals():
            collection.update_one(
                {"file_id": file_id},
                {"$set": {
                    "status": "failed",
                    "error": str(he),
                    "error_code": he.status_code,
                    "completed_at": datetime.utcnow()
                }}
            )
        raise he

    except Exception as e:
        # Update MongoDB with error status for general exceptions
        if 'file_id' in locals():
            collection.update_one(
                {"file_id": file_id},
                {"$set": {
                    "status": "failed",
                    "error": str(e),
                    "error_traceback": traceback.format_exc(),
                    "completed_at": datetime.utcnow()
                }}
            )
        logger.error(f"Error processing document: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Error processing document: {str(e)}")

    finally:
        # Clean up
        if 'pdf_document' in locals():
            pdf_document.close()
        await file.close()


def process_file_sync(file_path, file_id, openai_api_key):
    """
    Synchronously process a file from the given file path
    Args:
        file_path (str): Path to the file to process
        file_id (str): Unique identifier for the file
        openai_api_key (str): OpenAI API key for processing
    Returns:
        dict: Processing results
    """
    try:
        # Open the PDF document
        pdf_document = fitz.open(file_path)
        images_with_info = process_pdf(pdf_document)

        # Process images and collect results
        results = []
        annotated_image_paths = []
        
        for idx, (image, page_numbers, page_boundary) in enumerate(images_with_info):
            # Detect elements
            result_image, detected_elements = detect(image, page_numbers, page_boundary)

            # Save annotated image
            annotated_image_path = save_annotated_image(result_image, f"annotated_page_{idx}.png")
            annotated_image_paths.append(annotated_image_path)

            # Create and process chunks
            chunks = improved_intelligent_chunking_with_continuity(
                detected_elements, 
                SEGMENT_HIERARCHY, 
                max_chunk_size=10
            )

            processed_chunks = []
            for chunk_idx, chunk in enumerate(chunks):
                chunk_image = combine_elements_into_image(image, chunk["elements"])
                annotated_chunk_image, _ = detect(chunk_image)

                chunk["chunk_index"] = chunk_idx
                chunk["original_image"] = chunk_image
                chunk["annotated_image"] = annotated_chunk_image

                processed_chunk = process_chunk(chunk, openai_api_key)
                processed_chunks.append(processed_chunk)

            # Collect results for this page
            page_result = {
                "pages": page_numbers,
                "processed_chunks": processed_chunks,
                "detected_elements": detected_elements,
                "annotated_image_path": annotated_image_path
            }
            results.append(page_result)

        # Store the results in MongoDB
        collection.update_one(
            {"file_id": file_id},
            {
                "$set": {
                    "result": "Processing complete",
                    "chunks": results
                }
            },
            upsert=True
        )

        return results

    except Exception as e:
        logger.error(f"Error in process_file_sync: {str(e)}\n{traceback.format_exc()}")
        raise Exception(f"Error processing file: {str(e)}")
    finally:
        if 'pdf_document' in locals():
            pdf_document.close()

# Add the process_file_from_gcs function
def process_file_from_gcs(gcs_url, file_id, openai_api_key):
    # Parse the GCS URL to get bucket and blob names
    parsed_url = urlparse(gcs_url)
    path_parts = parsed_url.path.lstrip('/').split('/')
    bucket_name = path_parts[0]
    blob_name = '/'.join(path_parts[1:])
    
    # Get the blob
    bucket = storage_client.bucket(bucket_name)
    blob = bucket.blob(blob_name)
    
    # Download the file to a temporary location
    temp_dir = "temp_downloads"
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)
    temp_file_path = os.path.join(temp_dir, blob_name.split('/')[-1])
    blob.download_to_filename(temp_file_path)
    
    # Process the file
    all_results = process_file_sync(temp_file_path, file_id, openai_api_key)
    
    # Clean up the temporary file
    os.remove(temp_file_path)
    return all_results


# FastAPI Endpoints
@app.on_event("startup")
def startup_event():
    """
    Initialize the YOLO model on application startup
    """
    load_model()


@app.get("/response", response_model=ProcessResponse)
async def get_response(
    file_id: str = Query(...),
    gcs_url: Optional[str] = Query(None, description="GCS URL in the format 'gs://vision-bucket-ai/filename.pdf'")
):
    """
    Endpoint to retrieve the response stored in the database based on file_id.
    Optionally process a file from GCS if gcs_url is provided.
    
    Args:
        file_id: Unique identifier for the document
        gcs_url: Optional GCS URL in the format 'gs://vision-bucket-ai/filename.pdf'
    """
    try:
        # If gcs_url is provided, process the file from GCS
        if gcs_url:
            # Convert gs:// URL to https:// URL
            if gcs_url.startswith('gs://'):
                bucket_name = gcs_url.split('/')[2]
                blob_name = '/'.join(gcs_url.split('/')[3:])
                https_url = f"https://storage.googleapis.com/{bucket_name}/{blob_name}"
            else:
                raise HTTPException(status_code=400, detail="Invalid GCS URL format. Must start with 'gs://'")

            # Check if OpenAI configuration exists
            if not openai_config:
                raise HTTPException(status_code=400, detail="OpenAI API key not initialized")

            try:
                # Process the file from GCS
                results = process_file_from_gcs(https_url, file_id, openai_config.api_key)

                # Update MongoDB with the new results
                collection.update_one(
                    {"file_id": file_id},
                    {
                        "$set": {
                            "result": "Processing complete",
                            "chunks": results,
                            "gcs_url": https_url,
                            "updated_at": datetime.utcnow()
                        }
                    },
                    upsert=True
                )
            except Exception as e:
                logger.error(f"Error processing GCS file: {str(e)}")
                raise HTTPException(status_code=500, detail=f"Error processing GCS file: {str(e)}")

        # Query the MongoDB collection for the document with the given file_id
        response_data = collection.find_one({"file_id": file_id})
        if response_data is None:
            raise HTTPException(status_code=404, detail="Response not found")
        
        return ProcessResponse(
        
            result=response_data["result"],
            chunks=response_data["chunks"],
            gcs_url=response_data.get("gcs_url")
        )

    except HTTPException as he:
        # Re-raise HTTP exceptions
        raise he
    except Exception as e:
        # Log and raise any other unexpected errors
        logger.error(f"Unexpected error in get_response: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Unexpected error: {str(e)}")
    
@app.get("/list-templates")
async def list_templates():
    """
    List all available document templates
    """
    templates = await template_manager.list_templates()
    return templates

@app.get("/template/{template_id}")
async def get_template(template_id: str):
    """
    Get a specific template by ID
    """
    template = await template_manager.get_template(template_id)
    if template is None:
        raise HTTPException(status_code=404, detail="Template not found")
    return template

@app.post("/process-with-template", response_class=HTMLResponse)
async def process_with_template(
    template_id: str = Query(..., description="ID of the template to use"),
    test_file_type: str = Query(..., description="Type of test file to use", enum=["invoice", "balance_sheet", "bank_statement"])
):
    """
    Process a document using a specific template with enhanced extraction and summary generation.
    Only supports default test files.
    """
    try:
        # Check if OpenAI configuration is set
        if not openai_config:
            raise HTTPException(status_code=400, detail="OpenAI API key not initialized")

        # Handle test file selection
        if test_file_type not in DEFAULT_TEST_FILES:
            raise HTTPException(
                status_code=400,
                detail=f"Invalid test_file_type. Allowed values: {list(DEFAULT_TEST_FILES.keys())}"
            )

        file_path = DEFAULT_TEST_FILES[test_file_type]
        if not os.path.exists(file_path):
            raise HTTPException(
                status_code=404,
                detail=f"Default test file {test_file_type} not found. Please ensure default files are properly configured."
            )

        with open(file_path, 'rb') as f:
            content = f.read()
        
        filename = os.path.basename(file_path)
        file_type = "application/pdf"
        display_filename = filename

        # Generate a unique file_id for tracking
        file_id = f"{int(time.time())}_{filename}"

        # Get template configuration
        template = await template_manager.get_template(template_id)
        if not template:
            raise HTTPException(status_code=404, detail="Template not found")

        # Upload file to GCS
        file_name = f"{int(time.time())}_{filename}"
        blob = bucket.blob(file_name)
        blob.upload_from_string(content)
        gcs_url = f"https://storage.googleapis.com/{bucket_name}/{file_name}"

        # Create initial MongoDB document
        initial_doc = {
            "file_id": file_id,
            "gcs_url": gcs_url,
            "status": "processing",
            "result": "Processing started",
            "chunks": [],
            "file_name": filename,
            "template_id": template_id,
            "template_config": template,
            "created_at": datetime.utcnow(),
            "is_test_file": True,
            "test_file_type": test_file_type
        }
        collection.insert_one(initial_doc)

        # Process the PDF content
        pdf_document = fitz.open(stream=content, filetype="pdf")
        images_with_info = process_pdf(pdf_document)

        # Update MongoDB with processing started
        collection.update_one(
            {"file_id": file_id},
            {"$set": {
                "status": "processing_pages",
                "total_pages": len(images_with_info)
            }}
        )

        # Process images and collect results
        results = []
        annotated_image_paths = []
        
        for idx, (image, page_numbers, page_boundary) in enumerate(images_with_info):
            try:
                # Detect elements
                result_image, detected_elements = detect(image, page_numbers, page_boundary)

                # Save annotated image
                annotated_image_path = save_annotated_image(result_image, f"annotated_page_{idx}.png")
                annotated_image_paths.append(annotated_image_path)

                # Create and process chunks
                chunks = improved_intelligent_chunking_with_continuity(
                    detected_elements, 
                    SEGMENT_HIERARCHY, 
                    max_chunk_size=10
                )

                processed_chunks = []
                for chunk_idx, chunk in enumerate(chunks):
                    chunk_image = combine_elements_into_image(image, chunk["elements"])
                    annotated_chunk_image, _ = detect(chunk_image)

                    chunk["chunk_index"] = chunk_idx
                    chunk["original_image"] = chunk_image
                    chunk["annotated_image"] = annotated_chunk_image

                    processed_chunk = process_chunk(chunk, openai_config.api_key)
                    processed_chunks.append(processed_chunk)

                # Collect results for this page
                page_result = {
                    "pages": page_numbers,
                    "processed_chunks": processed_chunks,
                    "detected_elements": detected_elements,
                    "annotated_image_path": annotated_image_path
                }
                results.append(page_result)

                # Update MongoDB with page progress
                collection.update_one(
                    {"file_id": file_id},
                    {"$set": {
                        "pages_processed": idx + 1,
                        f"page_{idx}_status": "completed"
                    }}
                )

            except Exception as page_error:
                collection.update_one(
                    {"file_id": file_id},
                    {"$set": {
                        f"page_{idx}_status": "failed",
                        f"page_{idx}_error": str(page_error)
                    }}
                )
                logger.error(f"Error processing page {idx}: {str(page_error)}")
                continue

        overall_summary = generate_overall_summary(results, openai_config.api_key)

        # Update final results in MongoDB
        final_update = {
            "status": "completed",
            "result": "Processing complete",
            "chunks": results,
            "overall_summary": overall_summary,
            "completed_at": datetime.utcnow(),
            "annotated_image_paths": annotated_image_paths,
            "total_chunks_processed": sum(len(page["processed_chunks"]) for page in results)
        }

        collection.update_one(
            {"file_id": file_id},
            {"$set": final_update}
        )

        # Return the HTML response with file_id included
        return templates.TemplateResponse(
            "results.html",
            {
                "request": {"type": "http.request", "method": "POST"},
                "file_name": display_filename,
                "results": results,
                "overall_summary": overall_summary,
                "annotated_image_paths": annotated_image_paths,
                "file_id": file_id,
                "gcs_url": gcs_url,
                "template_id": template_id,
                "message": f"File ID for this document: {file_id}"
            }
        )

    except HTTPException as he:
        if 'file_id' in locals():
            collection.update_one(
                {"file_id": file_id},
                {"$set": {
                    "status": "failed",
                    "error": str(he),
                    "error_code": he.status_code,
                    "completed_at": datetime.utcnow()
                }}
            )
        raise he

    except Exception as e:
        if 'file_id' in locals():
            collection.update_one(
                {"file_id": file_id},
                {"$set": {
                    "status": "failed",
                    "error": str(e),
                    "error_traceback": traceback.format_exc(),
                    "completed_at": datetime.utcnow()
                }}
            )
        logger.error(f"Error processing document: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Error processing document: {str(e)}")

    finally:
        if 'pdf_document' in locals():
            pdf_document.close()


@app.delete("/delete-document/{file_id}")
async def delete_document(file_id: str):
    """
    Endpoint to delete a document and its associated data from MongoDB based on file_id.
    
    Args:
        file_id (str): The unique identifier of the document to delete
        
    Returns:
        dict: A message indicating the success or failure of the deletion
        
    Raises:
        HTTPException: If the document is not found or if there's an error during deletion
    """
    try:
        # Check if the document exists
        document = collection.find_one({"file_id": file_id})
        if not document:
            raise HTTPException(
                status_code=404,
                detail=f"Document with file_id {file_id} not found"
            )
        
        # Delete associated files if they exist
        try:
            # Delete converted file if it exists
            converted_filepath = os.path.join(CONVERTED_FILES_DIR, f"converted_{file_id}.pdf")
            if os.path.exists(converted_filepath):
                os.remove(converted_filepath)
            
            # Delete any associated annotated images
            static_dir = "static"
            if os.path.exists(static_dir):
                for filename in os.listdir(static_dir):
                    if file_id in filename:
                        file_path = os.path.join(static_dir, filename)
                        os.remove(file_path)
        except Exception as e:
            logging.warning(f"Error cleaning up files for {file_id}: {str(e)}")
        
        # Delete the document from MongoDB
        result = collection.delete_one({"file_id": file_id})
        
        if result.deleted_count == 1:
            logging.info(f"Successfully deleted document with file_id: {file_id}")
            return {
                "status": "success",
                "message": f"Document with file_id {file_id} has been successfully deleted",
                "deleted_count": result.deleted_count
            }
        else:
            raise HTTPException(
                status_code=500,
                detail="Document found but deletion failed"
            )
            
    except HTTPException as he:
        # Re-raise HTTP exceptions
        raise he
    except Exception as e:
        # Log the error and raise an HTTP exception
        logging.error(f"Error deleting document: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"Error deleting document: {str(e)}"
        )


# Health check endpoint
@app.get("/health")
async def health_check():
    """
    Basic health check endpoint
    """
    return {
        "status": "healthy", 
        "model_loaded": DETECTION_MODEL is not None
    }

# Catch-all exception handler
@app.exception_handler(Exception)
async def global_exception_handler(request, exc):
    """
    Global exception handler for unhandled exceptions
    """
    logging.error(f"Unhandled exception: {str(exc)}")
    return PlainTextResponse(
        status_code=500,
        content=f"An unexpected error occurred: {str(exc)}"
    )

# Example of how to run the FastAPI application
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
